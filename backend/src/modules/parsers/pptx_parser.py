from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.slide import Slide
from pptx.shapes.base import BaseShape
from pptx.table import Table

from src.config import settings
from src.services.tempfile_service import tempfile_service
from src.modules.parsers.abstract_parser import AbstractParser


class PPTXToMDParser(AbstractParser):
    def parse(self, filename: str, *args, **kwars) -> str:
        prs = Presentation(filename)
        markdown_content = []

        for slide_num, slide in enumerate(prs.slides, 1):
            title = self._get_slide_title(slide)
            if title:
                markdown_content.append(f"# {title} \n\n")
            else:
                markdown_content.append(f"# Слайд {slide_num}\n\n")

            for shape in slide.shapes:
                content = self._process_shape(shape, title)

                if content:
                    markdown_content.append(content)
                    markdown_content.append("\n")

        md_text = "\n".join(markdown_content)

        return self._preprocess_md(md_text)

    def _get_slide_title(self, slide: Slide) -> str | None:
        if slide.shapes.title:
            return slide.shapes.title.text.strip()
        return None

    def _process_shape(
        self,
        shape: BaseShape,
        slide_title: str | None,
    ) -> str | None:
        if shape.has_table:
            table_md = self._parse_table(shape.table)
            return table_md

        elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            image_md = self._parse_image(shape)
            if image_md:
                return image_md

        elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            group_content = []
            for grouped_shape in shape.shapes:
                content = self._process_shape(grouped_shape, slide_title)
                if content:
                    group_content.append(content)

            if group_content:
                return "\n".join(group_content)

        elif hasattr(shape, "text") and (shape_text := shape.text.strip()):
            if shape_text != slide_title:
                text_md = self._parse_text(shape)
                if text_md:
                    return text_md

        return None

    def _parse_table(self, table: Table) -> str:
        rows = []

        for i, row in enumerate(table.rows):
            cells = []
            for cell in row.cells:
                cell_text = cell.text.strip().replace("\n", "<br>")
                cells.append(cell_text)

            rows.append("| " + " | ".join(cells) + " |")

            if i == 0:
                separator = "| " + " | ".join(["---"] * len(cells)) + " |"
                rows.append(separator)

        return "\n".join(rows)

    def _parse_image(self, shape: BaseShape) -> str:
        image = shape.image

        # image_path = tempfile_service.save_file(image.blob, image.ext)

        # return f"![]({settings.DOMAIN}files/{image_path})"

        return self._image_to_md(image.blob, image.ext)

    def _parse_text(self, shape) -> str | None:
        if not hasattr(shape, "text_frame"):
            return shape.text.strip()

        text_parts = []

        for paragraph in shape.text_frame.paragraphs:
            text = paragraph.text.strip()
            if not text:
                continue

            level = paragraph.level

            if level > 0:
                indent = "  " * level
                text_parts.append(f"{indent}- {text}")
            else:
                if any(run.text.startswith(("•", "-", "*")) for run in paragraph.runs):
                    text_parts.append(f"- {text.lstrip('•-* ')}")
                else:
                    text_parts.append(text)

        return "\n".join(text_parts) if text_parts else None
