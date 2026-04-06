import fitz
from pptx import Presentation

from src.config import model_settings


class DocumentProcessor:
    def __init__(
        self,
        chunk_size: int = model_settings.CHUNK_SIZE,
        chunk_overlap: int = model_settings.CHUNK_OVERLAP,
    ):
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap

    def extract_text_from_pdf(self, path: str) -> str:
        if not fitz:
            raise RuntimeError("PyMuPDF (fitz) required. pip install pymupdf")
        text = []
        with fitz.open(path) as doc:
            for page in doc:
                text.append(page.get_text("text"))
        return "\n".join(text)

    def extract_text_from_pptx(self, path: str) -> str:
        if not Presentation:
            raise RuntimeError("python-pptx required. pip install python-pptx")
        prs = Presentation(path)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
        return "\n".join(text_runs)

    def chunk_text(self, text: str) -> list[str]:
        if not text or not text.strip():
            return []
        words = text.split()
        chunks = []
        start = 0
        while start < len(words):
            end = min(start + self.chunk_size, len(words))
            chunk = " ".join(words[start:end])
            chunks.append(chunk)
            start += self.chunk_size - self.chunk_overlap
        return chunks

    def process_document(
        self, document: str, metadata: dict | None = None
    ) -> list[dict]:
        path = metadata.get("source") if metadata else None
        ext = path.lower().split(".")[-1] if path else None

        if ext == "pdf":
            content = self.extract_text_from_pdf(path)
        elif ext in ("ppt", "pptx"):
            content = self.extract_text_from_pptx(path)
        else:
            content = document

        chunks = self.chunk_text(content)
        return [
            {"chunk_id": i, "text": ch, "metadata": metadata or {}}
            for i, ch in enumerate(chunks)
        ]
